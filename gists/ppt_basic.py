
from pptx import Presentation

def get_texts_from_file(input_file):
    
    # Instantiate variable to store the texts
    texts = []

    # Load the presentation
    prs = Presentation(input_file)

    # Fore each slide in tthe presentation
    for slide_number, slide in enumerate(prs.slides):  
        
        # For each shape in a slide
        for shape in slide.shapes:
    
            # Testing the "has_text_frame" parameter
            if shape.has_text_frame:
                
                # For each paragraph of the text_frame
                for paragraph in shape.text_frame.paragraphs:
                    
                    # Prepare new text
                    # texts.append(paragraph.text)
                    paragraph.text = ''

    
    prs.save(file_name.replace('.pptx', '_empty.pptx'))
    return texts

if __name__ == "__main__":
    
    file_name = '../ppt_translate.pptx'
    texts = get_texts_from_file(file_name)

    print('Texts has been selected.')
    print('Found {} paragraphs.'.format(len(texts)))