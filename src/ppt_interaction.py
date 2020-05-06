import os
from enum import IntEnum
from logzero import logger
from pptx import Presentation

# Defining Global vars 
import sys
config = sys.modules[__name__]
config.MODIF        = 0         # Indicates which modification to apply on text
config.CORPUS       = []        # Used to store texts of the PowerPoint
config.REPLACE      = {}        # Usedto replace part of the text (looking for caracters in a text)
config.TRANSLATION  = {}        # Used to replace full sentences   (looking for a text)


class TextModification(IntEnum):
    NO_MODIFICATION = 0
    EMPTY           = 1
    UPPER           = 2
    STORE_TEXT      = 3
    TRANSLATE       = 4
    REPLACE         = 5


def make_text_modification(text, modif=None):
    """Transform a text into another one.
       Key input : This the final function in ppt modification chain.
       Key input : Modification is made depenging on the Global Variables.
            - config.MODIF       : specify which transformation to realize
            - config.CORPUS      : needed list in case of storing text data
            - config.REPLACE     : needed dict for replacing part of strings
            - config.TRANSLATION : needed dict for translation purpose 

    PARAMETERS:   
        - text  : str - original text.
        - modif : int - specify the text modification to apply 
    RETURNS:
        - str : modified text.
    """
    if not modif:
        modif = config.MODIF
    
    if modif==TextModification.NO_MODIFICATION:
        return text

    if modif==TextModification.EMPTY:
        return ''
    
    if modif==TextModification.UPPER:
        return text.upper()

    if modif==TextModification.STORE_TEXT:
        config.CORPUS.append(text)
        return text
    
    if (modif==TextModification.TRANSLATE) & (text in config.TRANSLATION.keys()):
        return config.TRANSLATION[text]

    if modif==TextModification.REPLACE:
        for original_str, replace_str  in config.REPLACE.items():
            if original_str in text:
                text.replace(original_str, replace_str)
        return text

    logger.debug('No transformation on : {}'.format(text))
    return text


def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
    """ Given a paragraph 
    
    PARAMETERS:
        - paragraph : pptx.text.text._Paragraph - read paragraph with its original value and font parameters 
        - new_text : str - text to input into the paragraph
    
    INPLACE FUNCTION. NO OUTPUT.
    """
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text


def change_table_text(shape):
    """ Operate text modifications to a shape which contains/is a table.

    PARAMETERS:   
        - shape : pptx.shapes.shapetree.SlideShapes - original value.
    RETURNS:
        - pptx.shapes.shapetree.SlideShapes - modified values according to config.MODIF setting.
    """
    logger.warning('TABLE')

    # Get table information
    table = shape.table
    nb_rows = len(table.rows)
    nb_col = len(table.columns)

    # And update each cell
    for col in range(nb_col):
        for row in range(nb_rows):
            old_text = table.cell(row, col).text
            new_text = make_text_modification(old_text)
            table.cell(row, col).text = new_text
    return shape


def change_text_frame_text(shape):
    """
    Operate text modifications to a shape which contains/is a text_frame.

    PARAMETERS:   
        - shape : pptx.shapes.shapetree.SlideShapes - original value.
    RETURNS:
        - pptx.shapes.shapetree.SlideShapes - modified values according to config.MODIF setting.
    """
    # For each paragraph of the shape's text_frame
    text_frame = shape.text_frame
    for idx in range(len(text_frame.paragraphs)):
        
        # Prepare new text
        old_text = text_frame.paragraphs[idx].text
        new_text = make_text_modification(old_text)
        
        # Store it in presentation without modifying font parameters
        para = text_frame.paragraphs[idx]
        if len(para.runs)>0:
            replace_paragraph_text_retaining_initial_formatting(para, new_text)
        else:
            text_frame.paragraphs[idx].text = new_text

    return shape


def browse_shape(shape):
    """ PowerPoint Shapes might be of different kind.
        Each kind of shape has its own way to deal with text.
        This functions aims to root shapes to the right text extractor.

    PARAMETERS:   
        - shape : pptx.shapes.shapetree.SlideShapes - original value.
    RETURNS:
        - pptx.shapes.shapetree.SlideShapes - modified values according to config.MODIF setting.
    """
    
    # Text into text shape (basic case)
    if shape.has_text_frame:
        return change_text_frame_text(shape)

    # Text into tables (modification by cell)
    if shape.has_table:
        return change_table_text(shape)

    # Grouped shapes (using recursivity)
    if shape.shape_type==6:
        logger.warning('GROUPED')
        return browse_slide(shape)

    # If other cases specify it. Should not happen.
    # logger.warn('Unknown type of shape.')
    # logger.warn(shape.shape_type)    
    return shape


def browse_slide(slide):
    """ All PowerPoint slides are composed of shapes. We're here dealing with each.
        Applying modifications, depending on global variable config.MODIF.

    PARAMETERS:   
        - slide : pptx.slide.Slides (or pptx.shapes.shapetree.GroupShapes) - original values.
    RETURNS:
        - pptx.slide.Slides (or pptx.shapes.shapetree.GroupShapes) - modified values according to config.MODIF setting.
    """
    for shape in slide.shapes:
        shape = browse_shape(shape)
    return slide

    

def browse_file(input_file, output_file=None, text_modif=0, open_file=False): 
    """Browse a PowerPoint file and apply modification to its text.

    PARAMETERS : 
        - inpute_file : str - path of the fle to browse
        - output_file : str - path to save modified PowerPoint. (Default is None : file is not saved.)
        - text_modif : integer specifying the transformation to apply. See TextModification class for more.
        - open_file : boolean - if output_file is set, this will dirrectly open the PowerPoint file after browsing ad saving.
    """

    # Browsing file
    config.MODIF = text_modif
    logger.error(config.MODIF)
    prs = Presentation(input_file)
    for slide_number, slide in enumerate(prs.slides):  
        logger.info('> NEW Slide : {}'.format(slide_number + 1))
        slide = browse_slide(slide)

    # Displaying information to user
    logger.info('File Browsed : {}'.format(output_file))
    if output_file:
        prs.save(output_file)
        logger.info('File Saved   : {}'.format(output_file))
        if open_file:
            logger.info('File Opening : {}'.format(output_file))
            os.system('open {}'.format(output_file))

